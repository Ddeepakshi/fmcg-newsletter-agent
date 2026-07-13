"""Deliverable writers: CSV, JSON, Word, Excel, PowerPoint (spec section 11).

Each writer takes the final structured deal records (+ newsletter text where
relevant) and produces one deliverable file. Kept separate from pipeline.py
so any single output format can be regenerated without re-running the
pipeline.

Word and PowerPoint are styled as an actual newsletter/deck (masthead,
accent colors, real clickable hyperlinks, one slide per deal) rather than
plain unstyled text, since these are the two formats a business user is
meant to actually read rather than just data-dump into.
"""
import json
import logging
import re
from datetime import datetime, timezone

import pandas as pd
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.opc.constants import RELATIONSHIP_TYPE
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Pt, RGBColor
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt as PptPt

from src.schema import CSV_COLUMN_ORDER

logger = logging.getLogger(__name__)

_BOLD_RE = re.compile(r"\*\*(.+?)\*\*")
_LINK_RE = re.compile(r"\[([^\]]+)\]\(([^)]+)\)")
_INLINE_TOKEN_RE = re.compile(r"\*\*(?P<bold>.+?)\*\*|\[(?P<link_text>[^\]]+)\]\((?P<link_url>[^)]+)\)")

# A single professional accent palette shared by both documents.
ACCENT = RGBColor(0x1B, 0x5E, 0x20)       # deep FMCG green
ACCENT_LIGHT = RGBColor(0x4C, 0xAF, 0x50)
MUTED = RGBColor(0x75, 0x75, 0x75)
TEXT_DARK = RGBColor(0x21, 0x21, 0x21)
LINK_BLUE = RGBColor(0x11, 0x55, 0xCC)

PPT_ACCENT = PptRGBColor(0x1B, 0x5E, 0x20)
PPT_ACCENT_LIGHT = PptRGBColor(0x4C, 0xAF, 0x50)
PPT_MUTED = PptRGBColor(0x75, 0x75, 0x75)
PPT_TEXT_DARK = PptRGBColor(0x21, 0x21, 0x21)
PPT_WHITE = PptRGBColor(0xFF, 0xFF, 0xFF)


def _markdown_to_plain(line: str) -> str:
    line = _LINK_RE.sub(lambda m: f"{m.group(1)} ({m.group(2)})", line)
    line = _BOLD_RE.sub(lambda m: m.group(1), line)
    return line


def _parse_inline_markdown(text: str) -> list:
    """Tokenizes `**bold**` and `[text](url)` into (text, bold, url) triples
    so callers can render rich runs (bold/hyperlink) instead of flattening
    to plain text.
    """
    tokens = []
    pos = 0
    for m in _INLINE_TOKEN_RE.finditer(text):
        if m.start() > pos:
            tokens.append((text[pos:m.start()], False, None))
        if m.group("bold") is not None:
            tokens.append((m.group("bold"), True, None))
        else:
            tokens.append((m.group("link_text"), False, m.group("link_url")))
        pos = m.end()
    if pos < len(text):
        tokens.append((text[pos:], False, None))
    return tokens or [(text, False, None)]


def to_dataframe(records: list) -> pd.DataFrame:
    df = pd.DataFrame(records)
    for col in CSV_COLUMN_ORDER:
        if col not in df.columns:
            df[col] = None
    return df[CSV_COLUMN_ORDER]


def write_csv(records: list, path: str) -> None:
    to_dataframe(records).to_csv(path, index=False)
    logger.info("Wrote CSV: %s (%d rows)", path, len(records))


def write_json(records: list, path: str) -> None:
    with open(path, "w", encoding="utf-8") as f:
        json.dump(records, f, indent=2, ensure_ascii=False)
    logger.info("Wrote JSON: %s (%d rows)", path, len(records))


def write_xlsx(records: list, path: str, newsletter_markdown: str = None) -> None:
    """Writes the structured deal table, plus a second "Newsletter" sheet with
    the actual newsletter text (one line per row) when provided — so the
    Excel deliverable carries the newsletter content itself, not just the
    raw data table behind it.
    """
    df = to_dataframe(records)
    with pd.ExcelWriter(path, engine="openpyxl") as writer:
        df.to_excel(writer, index=False, sheet_name="Deals")
        worksheet = writer.sheets["Deals"]
        for i, col in enumerate(df.columns, start=1):
            max_len = max([len(str(col))] + [len(str(v)) for v in df[col].fillna("")])
            worksheet.column_dimensions[get_column_letter(i)].width = min(max_len + 2, 60)

        if newsletter_markdown:
            plain_lines = [_markdown_to_plain(line) for line in newsletter_markdown.splitlines()]
            newsletter_df = pd.DataFrame({"Newsletter": plain_lines})
            newsletter_df.to_excel(writer, index=False, sheet_name="Newsletter")
            writer.sheets["Newsletter"].column_dimensions["A"].width = 100

    logger.info(
        "Wrote Excel: %s (%d deal rows%s)",
        path, len(records), ", + Newsletter sheet" if newsletter_markdown else "",
    )


def _add_docx_hyperlink_run(paragraph, url: str, text: str, bold: bool = False):
    """python-docx has no first-class hyperlink API — this is the standard
    low-level OOXML workaround for a real, clickable link run.
    """
    part = paragraph.part
    r_id = part.relate_to(url, RELATIONSHIP_TYPE.HYPERLINK, is_external=True)

    hyperlink = OxmlElement("w:hyperlink")
    hyperlink.set(qn("r:id"), r_id)

    run = OxmlElement("w:r")
    rPr = OxmlElement("w:rPr")
    color = OxmlElement("w:color")
    color.set(qn("w:val"), "1155CC")
    rPr.append(color)
    underline = OxmlElement("w:u")
    underline.set(qn("w:val"), "single")
    rPr.append(underline)
    if bold:
        rPr.append(OxmlElement("w:b"))
    run.append(rPr)

    t = OxmlElement("w:t")
    t.text = text
    run.append(t)
    hyperlink.append(run)
    paragraph._p.append(hyperlink)


def _add_docx_horizontal_rule(paragraph):
    pPr = paragraph._p.get_or_add_pPr()
    pBdr = OxmlElement("w:pBdr")
    bottom = OxmlElement("w:bottom")
    bottom.set(qn("w:val"), "single")
    bottom.set(qn("w:sz"), "8")
    bottom.set(qn("w:space"), "1")
    bottom.set(qn("w:color"), "1B5E20")
    pBdr.append(bottom)
    pPr.append(pBdr)


def _add_docx_rich_paragraph(doc, text: str, style: str = None, base_size: Pt = None):
    """Renders **bold**/[text](url) tokens as real bold runs / hyperlinks
    instead of flattening to plain text.
    """
    p = doc.add_paragraph(style=style) if style else doc.add_paragraph()
    for token_text, bold, url in _parse_inline_markdown(text):
        if not token_text:
            continue
        if url:
            _add_docx_hyperlink_run(p, url, token_text, bold=bold)
        else:
            run = p.add_run(token_text)
            run.bold = bold
            if base_size:
                run.font.size = base_size
    return p


def write_docx(newsletter_markdown: str, path: str) -> None:
    doc = Document()

    normal = doc.styles["Normal"]
    normal.font.name = "Calibri"
    normal.font.size = Pt(11)
    normal.font.color.rgb = TEXT_DARK

    first_heading = True
    for raw_line in newsletter_markdown.splitlines():
        line = raw_line.rstrip()
        if not line.strip():
            doc.add_paragraph("")
            continue

        if line.startswith("# "):
            heading_text = _markdown_to_plain(line[2:])
            h = doc.add_heading(level=1)
            h.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run = h.add_run(heading_text)
            run.font.color.rgb = ACCENT
            run.font.size = Pt(26)

            if first_heading:
                subtitle = doc.add_paragraph()
                subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
                sub_run = subtitle.add_run(
                    f"Recent M&A, Investment & Funding Activity  ·  Generated "
                    f"{datetime.now(timezone.utc).strftime('%B %d, %Y')}"
                )
                sub_run.italic = True
                sub_run.font.size = Pt(11)
                sub_run.font.color.rgb = MUTED
                _add_docx_horizontal_rule(subtitle)
                first_heading = False

        elif line.startswith("## "):
            h = doc.add_heading(level=2)
            run = h.add_run(_markdown_to_plain(line[3:]))
            run.font.color.rgb = ACCENT
            run.font.size = Pt(16)
        elif line.startswith("### "):
            h = doc.add_heading(level=3)
            run = h.add_run(_markdown_to_plain(line[4:]))
            run.font.color.rgb = ACCENT_LIGHT
            run.font.size = Pt(13)
        elif line.strip().startswith("- ") or line.strip().startswith("* "):
            _add_docx_rich_paragraph(doc, line.strip()[2:], style="List Bullet")
        elif line.strip().startswith("*") and line.strip().endswith("*") and len(line.strip()) > 1:
            p = doc.add_paragraph()
            run = p.add_run(_markdown_to_plain(line.strip().strip("*")))
            run.italic = True
            run.font.size = Pt(9)
            run.font.color.rgb = MUTED
        else:
            _add_docx_rich_paragraph(doc, line)

    footer = doc.add_paragraph()
    _add_docx_horizontal_rule(footer)
    footer_run = footer.add_run(
        "Generated by the FMCG Deal Intelligence Agent. Every entry links back to its original source."
    )
    footer_run.italic = True
    footer_run.font.size = Pt(9)
    footer_run.font.color.rgb = MUTED

    doc.save(path)
    logger.info("Wrote Word doc: %s", path)


def _ppt_accent_bar(slide, prs, height=Inches(0.12)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), prs.slide_width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = PPT_ACCENT
    bar.line.fill.background()
    bar.shadow.inherit = False
    return bar


def _ppt_textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    return tf


def _ppt_add_run(paragraph, text, size, color, bold=False, italic=False, url=None):
    run = paragraph.add_run()
    run.text = text
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    if url:
        run.hyperlink.address = url
    return run


def write_pptx(records: list, path: str, top_n: int = 10, title: str = "FMCG Deal Intelligence") -> None:
    """One title slide, one slide per top deal (by credibility), one closing
    summary slide — a real slide a business audience can present from,
    rather than a single slide with every deal crammed into bullet text.
    """
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    ranked = sorted(records, key=lambda r: r.get("credibility_score") or 0, reverse=True)[:top_n]
    generated_on = datetime.now(timezone.utc).strftime("%B %d, %Y")

    # --- Title slide ---
    slide = prs.slides.add_slide(blank_layout)
    _ppt_accent_bar(slide, prs, height=Inches(0.25))
    tf = _ppt_textbox(slide, Inches(0.7), Inches(2.3), prs.slide_width - Inches(1.4), Inches(1.5))
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    _ppt_add_run(tf.paragraphs[0], title, PptPt(40), PPT_TEXT_DARK, bold=True)
    sub = tf.add_paragraph()
    sub.alignment = PP_ALIGN.CENTER
    _ppt_add_run(sub, "Recent M&A, Investment & Funding Activity", PptPt(20), PPT_ACCENT)
    date_p = tf.add_paragraph()
    date_p.alignment = PP_ALIGN.CENTER
    _ppt_add_run(date_p, f"Generated {generated_on}", PptPt(13), PPT_MUTED, italic=True)

    if not ranked:
        slide = prs.slides.add_slide(blank_layout)
        _ppt_accent_bar(slide, prs)
        tf = _ppt_textbox(slide, Inches(0.7), Inches(2.5), prs.slide_width - Inches(1.4), Inches(1.5))
        _ppt_add_run(tf.paragraphs[0], "No qualifying deals found in this window.", PptPt(24), PPT_MUTED)

    # --- One slide per deal ---
    for i, r in enumerate(ranked, 1):
        slide = prs.slides.add_slide(blank_layout)
        _ppt_accent_bar(slide, prs)

        badge_tf = _ppt_textbox(slide, Inches(0.5), Inches(0.45), Inches(1.0), Inches(0.5))
        _ppt_add_run(badge_tf.paragraphs[0], f"{i:02d}", PptPt(18), PPT_ACCENT_LIGHT, bold=True)

        title_tf = _ppt_textbox(slide, Inches(0.5), Inches(1.0), prs.slide_width - Inches(1.0), Inches(1.5))
        _ppt_add_run(title_tf.paragraphs[0], r.get("title", "Untitled"), PptPt(26), PPT_TEXT_DARK, bold=True)

        highlight_parts = [p for p in (r.get("acquirer"), r.get("target")) if p]
        highlight = " → ".join(highlight_parts) if highlight_parts else None

        detail_bits = []
        if r.get("deal_type"):
            detail_bits.append(r["deal_type"])
        if r.get("deal_value"):
            detail_bits.append(f"{r.get('currency') or ''} {r['deal_value']}".strip())
        detail_line = " · ".join(detail_bits)

        info_tf = _ppt_textbox(slide, Inches(0.5), Inches(2.6), prs.slide_width - Inches(1.0), Inches(1.3))
        first = True
        if highlight:
            _ppt_add_run(info_tf.paragraphs[0], highlight, PptPt(20), PPT_ACCENT, bold=True)
            first = False
        if detail_line:
            p = info_tf.paragraphs[0] if first else info_tf.add_paragraph()
            _ppt_add_run(p, detail_line, PptPt(16), PPT_TEXT_DARK)
            first = False
        if r.get("region"):
            p = info_tf.paragraphs[0] if first else info_tf.add_paragraph()
            _ppt_add_run(p, f"Region: {r['region']}", PptPt(13), PPT_MUTED)

        footer_tf = _ppt_textbox(slide, Inches(0.5), prs.slide_height - Inches(0.9),
                                  prs.slide_width - Inches(1.0), Inches(0.6))
        source = r.get("source") or "Unknown source"
        url = r.get("url")
        credibility = r.get("credibility_score")
        cred_text = f"  ·  Credibility {credibility:.2f}" if credibility is not None else ""
        _ppt_add_run(footer_tf.paragraphs[0], f"Source: {source}{cred_text}", PptPt(12), PPT_MUTED, url=url)

    # --- Closing summary slide ---
    slide = prs.slides.add_slide(blank_layout)
    _ppt_accent_bar(slide, prs)
    tf = _ppt_textbox(slide, Inches(0.7), Inches(1.0), prs.slide_width - Inches(1.4), Inches(0.8))
    _ppt_add_run(tf.paragraphs[0], "Summary", PptPt(32), PPT_ACCENT, bold=True)

    stats_tf = _ppt_textbox(slide, Inches(0.7), Inches(2.0), prs.slide_width - Inches(1.4), Inches(2.5))
    avg_credibility = (
        sum(r.get("credibility_score") or 0 for r in ranked) / len(ranked) if ranked else 0
    )
    stats = [
        f"{len(ranked)} deal(s) featured in this edition",
        f"Average credibility score: {avg_credibility:.2f}",
        f"Generated {generated_on}",
        "Full dataset (with every source link) available in the accompanying CSV/JSON.",
    ]
    for i, stat in enumerate(stats):
        p = stats_tf.paragraphs[0] if i == 0 else stats_tf.add_paragraph()
        _ppt_add_run(p, f"•  {stat}", PptPt(18), PPT_TEXT_DARK)

    prs.save(path)
    logger.info("Wrote PowerPoint: %s (%d deal slides + title + summary)", path, len(ranked))


def write_all_outputs(records: list, newsletter_markdown: str, output_dir) -> dict:
    """Writes csv/json/docx/xlsx/pptx into `output_dir`, returns their paths."""
    from pathlib import Path

    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    paths = {
        "csv": output_dir / "deals.csv",
        "json": output_dir / "deals.json",
        "docx": output_dir / "newsletter.docx",
        "xlsx": output_dir / "deals.xlsx",
        "pptx": output_dir / "newsletter.pptx",
    }

    write_csv(records, str(paths["csv"]))
    write_json(records, str(paths["json"]))
    write_docx(newsletter_markdown, str(paths["docx"]))
    write_xlsx(records, str(paths["xlsx"]), newsletter_markdown=newsletter_markdown)
    write_pptx(records, str(paths["pptx"]))

    return {k: str(v) for k, v in paths.items()}


if __name__ == "__main__":
    import os
    import tempfile

    sample_records = [
        {"title": "ITC to acquire Prasuma Foods for $150 million", "snippet": "Deal expected to close in Q3",
         "source": "Economic Times", "source_tier": 1, "published_date": "2026-07-11T00:00:00+00:00",
         "url": "https://a.com/1", "region": "India", "is_duplicate_of": None, "credibility_score": 0.91,
         "acquirer": "ITC", "target": "Prasuma Foods", "deal_type": "Acquisition", "deal_value": 150,
         "currency": "USD", "relevance_flag": "Relevant"},
    ]
    sample_newsletter = (
        "# FMCG Deal Intelligence Newsletter\n\n## Recent Deal Activity\n\n"
        "- **ITC to acquire Prasuma Foods for $150 million**\n  Source: [Economic Times](https://a.com/1)\n"
    )

    with tempfile.TemporaryDirectory() as tmp_dir:
        paths = write_all_outputs(sample_records, sample_newsletter, tmp_dir)
        for fmt, path in paths.items():
            size = os.path.getsize(path)
            print(f"{fmt}: {path} ({size} bytes)")
            assert size > 0, f"{fmt} output is empty"

        from openpyxl import load_workbook

        wb = load_workbook(paths["xlsx"])
        assert "Newsletter" in wb.sheetnames, "xlsx should carry a Newsletter sheet, not just raw deal data"

    print("All output writers produced non-empty files.")
