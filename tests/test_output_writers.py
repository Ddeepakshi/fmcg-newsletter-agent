"""Tests for src/output_writers.py."""
import csv
import json

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

from src.output_writers import write_all_outputs
from src.schema import CSV_COLUMN_ORDER

SAMPLE_RECORDS = [
    {"title": "ITC to acquire Prasuma Foods for $150 million", "snippet": "Deal expected to close in Q3",
     "source": "Economic Times", "source_tier": 1, "published_date": "2026-07-11T00:00:00+00:00",
     "url": "https://a.com/1", "region": "India", "is_duplicate_of": None, "credibility_score": 0.91,
     "acquirer": "ITC", "target": "Prasuma Foods", "deal_type": "Acquisition", "deal_value": 150,
     "currency": "USD", "relevance_flag": "Relevant"},
]
SAMPLE_NEWSLETTER = (
    "# FMCG Deal Intelligence Newsletter\n\n## Recent Deal Activity\n\n"
    "- **ITC to acquire Prasuma Foods for $150 million**\n  Source: [Economic Times](https://a.com/1)\n"
)


def test_write_all_outputs_produces_all_five_formats(tmp_path):
    paths = write_all_outputs(SAMPLE_RECORDS, SAMPLE_NEWSLETTER, tmp_path)
    assert set(paths.keys()) == {"csv", "json", "docx", "xlsx", "pptx"}
    for path in paths.values():
        assert (tmp_path / path.split("/")[-1]).stat().st_size > 0


def test_csv_output_has_schema_columns(tmp_path):
    paths = write_all_outputs(SAMPLE_RECORDS, SAMPLE_NEWSLETTER, tmp_path)
    with open(paths["csv"], newline="", encoding="utf-8") as f:
        reader = csv.DictReader(f)
        assert reader.fieldnames == CSV_COLUMN_ORDER
        rows = list(reader)
    assert len(rows) == 1
    assert rows[0]["acquirer"] == "ITC"


def test_json_output_round_trips(tmp_path):
    paths = write_all_outputs(SAMPLE_RECORDS, SAMPLE_NEWSLETTER, tmp_path)
    with open(paths["json"], encoding="utf-8") as f:
        data = json.load(f)
    assert data == SAMPLE_RECORDS


def test_docx_output_is_valid_and_contains_heading(tmp_path):
    paths = write_all_outputs(SAMPLE_RECORDS, SAMPLE_NEWSLETTER, tmp_path)
    doc = Document(paths["docx"])
    text = "\n".join(p.text for p in doc.paragraphs)
    assert "FMCG Deal Intelligence Newsletter" in text
    assert "ITC to acquire Prasuma Foods" in text


def test_docx_output_has_a_real_clickable_hyperlink(tmp_path):
    # Guardrail against regressing to plain "text (url)" — sources should be
    # actual clickable hyperlinks, not flattened text.
    paths = write_all_outputs(SAMPLE_RECORDS, SAMPLE_NEWSLETTER, tmp_path)
    doc = Document(paths["docx"])
    hyperlink_targets = [rel.target_ref for rel in doc.part.rels.values() if "hyperlink" in rel.reltype]
    assert "https://a.com/1" in hyperlink_targets


def test_xlsx_output_is_valid(tmp_path):
    paths = write_all_outputs(SAMPLE_RECORDS, SAMPLE_NEWSLETTER, tmp_path)
    wb = load_workbook(paths["xlsx"])
    ws = wb["Deals"]
    assert ws.max_row == 2  # header + 1 data row
    assert ws.max_column == len(CSV_COLUMN_ORDER)


def test_xlsx_output_includes_newsletter_sheet(tmp_path):
    # The Excel deliverable should carry the newsletter content itself, not
    # just the raw deal table behind it.
    paths = write_all_outputs(SAMPLE_RECORDS, SAMPLE_NEWSLETTER, tmp_path)
    wb = load_workbook(paths["xlsx"])
    assert "Newsletter" in wb.sheetnames
    ws = wb["Newsletter"]
    all_text = "\n".join(str(row[0].value) for row in ws.iter_rows(min_row=2))
    assert "FMCG Deal Intelligence Newsletter" in all_text
    assert "ITC to acquire Prasuma Foods" in all_text


def test_pptx_output_is_valid(tmp_path):
    # Title slide + one slide per deal + a closing summary slide.
    paths = write_all_outputs(SAMPLE_RECORDS, SAMPLE_NEWSLETTER, tmp_path)
    prs = Presentation(paths["pptx"])
    assert len(prs.slides) == len(SAMPLE_RECORDS) + 2


def test_pptx_deal_slide_has_structured_content_and_hyperlink(tmp_path):
    paths = write_all_outputs(SAMPLE_RECORDS, SAMPLE_NEWSLETTER, tmp_path)
    prs = Presentation(paths["pptx"])
    deal_slide = prs.slides[1]
    all_text = "\n".join(sh.text_frame.text for sh in deal_slide.shapes if sh.has_text_frame)
    assert "ITC to acquire Prasuma Foods" in all_text
    assert "ITC" in all_text and "Prasuma Foods" in all_text
    assert "Economic Times" in all_text

    hyperlink_runs = [
        run for sh in deal_slide.shapes if sh.has_text_frame
        for p in sh.text_frame.paragraphs for run in p.runs
        if run.hyperlink.address
    ]
    assert any(run.hyperlink.address == "https://a.com/1" for run in hyperlink_runs)


def test_pptx_handles_no_qualifying_deals(tmp_path):
    paths = write_all_outputs([], "# FMCG Deal Intelligence Newsletter\n\nNo qualifying deals found.", tmp_path)
    prs = Presentation(paths["pptx"])
    assert len(prs.slides) == 3  # title + "no deals" notice + summary
